import logging
from typing import List, Optional, Dict, Any
from pathlib import Path
import mimetypes
from abc import ABC, abstractmethod
import functools
import concurrent.futures
from io import StringIO

# Import existing utilities
from .pdf_utils import PDFFileLoader
from .text_utils import TextFileLoader

class FileProcessorBase(ABC):
    """Abstract base class for file processors."""
    
    @abstractmethod
    def load_documents(self) -> List[str]:
        """Extract text content from the file."""
        pass
    
    @abstractmethod
    def get_metadata(self) -> Dict[str, Any]:
        """Get metadata from the file."""
        pass

class PDFProcessor(FileProcessorBase):
    """PDF file processor using existing PDFFileLoader."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.loader = PDFFileLoader(file_path)
    
    def load_documents(self) -> List[str]:
        return self.loader.load_documents()
    
    def get_metadata(self) -> Dict[str, Any]:
        metadata = self.loader.get_metadata()
        metadata.update({
            "file_type": "pdf",
            "page_count": self.loader.get_page_count()
        })
        return metadata

class TextProcessor(FileProcessorBase):
    """Plain text file processor."""
    
    def __init__(self, file_path: str, encoding: str = "utf-8"):
        self.file_path = file_path
        self.encoding = encoding
    
    @functools.lru_cache(maxsize=1)
    def load_documents(self) -> List[str]:
        try:
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                content = f.read().strip()
                return [content] if content else []
        except Exception as e:
            logging.error(f"Failed to load text file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        file_path = Path(self.file_path)
        return {
            "file_type": "text",
            "encoding": self.encoding,
            "file_size": file_path.stat().st_size,
            "filename": file_path.name
        }

class MarkdownProcessor(FileProcessorBase):
    """Markdown file processor."""
    
    def __init__(self, file_path: str, encoding: str = "utf-8"):
        self.file_path = file_path
        self.encoding = encoding
    
    @functools.lru_cache(maxsize=1)
    def load_documents(self) -> List[str]:
        try:
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                content = f.read().strip()
                return [content] if content else []
        except Exception as e:
            logging.error(f"Failed to load markdown file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        file_path = Path(self.file_path)
        return {
            "file_type": "markdown",
            "encoding": self.encoding,
            "file_size": file_path.stat().st_size,
            "filename": file_path.name
        }

class CSVProcessor(FileProcessorBase):
    """CSV file processor that treats each row as a document."""
    
    def __init__(self, file_path: str, encoding: str = "utf-8"):
        self.file_path = file_path
        self.encoding = encoding
    
    @functools.lru_cache(maxsize=1)
    def load_documents(self) -> List[str]:
        try:
            import csv
            documents = []
            buffer = StringIO()
            
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                # Try to detect delimiter
                sample = f.read(1024)
                f.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.DictReader(f, delimiter=delimiter)
                
                # Process rows in batches
                batch = []
                for i, row in enumerate(reader):
                    # Convert row to text representation
                    row_text = "\n".join([f"{key}: {value}" for key, value in row.items() if value])
                    if row_text.strip():
                        batch.append(f"Row {i+1}:\n{row_text}")
                    
                    # Process batch when it reaches size 100
                    if len(batch) >= 100:
                        documents.extend(batch)
                        batch = []
                
                # Add remaining rows
                if batch:
                    documents.extend(batch)
            
            return documents
        except Exception as e:
            logging.error(f"Failed to load CSV file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        file_path = Path(self.file_path)
        return {
            "file_type": "csv",
            "encoding": self.encoding,
            "file_size": file_path.stat().st_size,
            "filename": file_path.name
        }

class DOCXProcessor(FileProcessorBase):
    """DOCX file processor for Microsoft Word documents."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    @functools.lru_cache(maxsize=1)
    def load_documents(self) -> List[str]:
        try:
            from docx import Document
            
            doc = Document(self.file_path)
            paragraphs = []
            
            # Process paragraphs in parallel
            def process_paragraph(para):
                text = para.text.strip()
                return text if text else None
            
            with concurrent.futures.ThreadPoolExecutor() as executor:
                processed_paragraphs = list(executor.map(process_paragraph, doc.paragraphs))
            
            # Filter out None values and join paragraphs
            paragraphs = [p for p in processed_paragraphs if p]
            content = "\n\n".join(paragraphs)
            return [content] if content else []
            
        except ImportError:
            raise ImportError("python-docx package is required to process DOCX files. Install with: pip install python-docx")
        except Exception as e:
            logging.error(f"Failed to load DOCX file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        try:
            from docx import Document
            doc = Document(self.file_path)
            
            file_path = Path(self.file_path)
            metadata = {
                "file_type": "docx",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name,
                "paragraph_count": len(doc.paragraphs)
            }
            
            # Add document properties if available
            if doc.core_properties:
                core_props = doc.core_properties
                if core_props.title:
                    metadata["title"] = core_props.title
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.created:
                    metadata["created"] = core_props.created.isoformat()
                if core_props.modified:
                    metadata["modified"] = core_props.modified.isoformat()
            
            return metadata
            
        except ImportError:
            file_path = Path(self.file_path)
            return {
                "file_type": "docx",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }
        except Exception as e:
            logging.warning(f"Failed to extract DOCX metadata: {e}")
            file_path = Path(self.file_path)
            return {
                "file_type": "docx",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }

class ExcelProcessor(FileProcessorBase):
    """Excel file processor for .xlsx and .xls files."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load_documents(self) -> List[str]:
        try:
            import pandas as pd
            
            # Read all sheets from the Excel file
            documents = []
            with pd.ExcelFile(self.file_path) as excel_file:
                for sheet_name in excel_file.sheet_names:
                    try:
                        df = pd.read_excel(excel_file, sheet_name=sheet_name)
                        
                        # Skip empty sheets
                        if df.empty:
                            continue
                            
                        # Convert to markdown table format
                        markdown_table = df.to_markdown(index=False)
                        
                        # Create a document for this sheet
                        sheet_content = f"Sheet: {sheet_name}\n\n{markdown_table}"
                        documents.append(sheet_content)
                        
                    except Exception as e:
                        logging.warning(f"Failed to process sheet '{sheet_name}': {e}")
                        continue
            
            return documents
            
        except ImportError:
            raise ImportError("pandas and openpyxl packages are required to process Excel files. Install with: pip install pandas openpyxl")
        except Exception as e:
            logging.error(f"Failed to load Excel file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        try:
            import pandas as pd
            
            file_path = Path(self.file_path)
            metadata = {
                "file_type": "excel",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }
            
            # Add sheet information
            try:
                with pd.ExcelFile(self.file_path) as excel_file:
                    metadata["sheet_count"] = len(excel_file.sheet_names)
                    metadata["sheet_names"] = excel_file.sheet_names
            except Exception as e:
                logging.warning(f"Failed to extract Excel metadata: {e}")
            
            return metadata
            
        except ImportError:
            file_path = Path(self.file_path)
            return {
                "file_type": "excel",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }

class PowerPointProcessor(FileProcessorBase):
    """PowerPoint file processor for .pptx files."""
    
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load_documents(self) -> List[str]:
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            documents = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                slide_content.append(f"Slide {i + 1}:")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                # Extract notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    slide_content.append(f"Notes: {slide.notes_slide.notes_text_frame.text.strip()}")
                
                if len(slide_content) > 1:  # More than just the slide number
                    documents.append("\n\n".join(slide_content))
            
            return documents
            
        except ImportError:
            raise ImportError("python-pptx package is required to process PowerPoint files. Install with: pip install python-pptx")
        except Exception as e:
            logging.error(f"Failed to load PowerPoint file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            file_path = Path(self.file_path)
            
            metadata = {
                "file_type": "powerpoint",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name,
                "slide_count": len(prs.slides)
            }
            
            # Add document properties if available
            if prs.core_properties:
                core_props = prs.core_properties
                if core_props.title:
                    metadata["title"] = core_props.title
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.created:
                    metadata["created"] = core_props.created.isoformat()
                if core_props.modified:
                    metadata["modified"] = core_props.modified.isoformat()
            
            return metadata
            
        except ImportError:
            file_path = Path(self.file_path)
            return {
                "file_type": "powerpoint",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }
        except Exception as e:
            logging.warning(f"Failed to extract PowerPoint metadata: {e}")
            file_path = Path(self.file_path)
            return {
                "file_type": "powerpoint",
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }

class HTMLProcessor(FileProcessorBase):
    """HTML file processor for .html and .htm files."""
    
    def __init__(self, file_path: str, encoding: str = "utf-8"):
        self.file_path = file_path
        self.encoding = encoding
    
    def load_documents(self) -> List[str]:
        try:
            from bs4 import BeautifulSoup
            
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                html_content = f.read()
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extract text content
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return [text] if text.strip() else []
            
        except ImportError:
            raise ImportError("beautifulsoup4 package is required to process HTML files. Install with: pip install beautifulsoup4")
        except Exception as e:
            logging.error(f"Failed to load HTML file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        try:
            from bs4 import BeautifulSoup
            
            file_path = Path(self.file_path)
            metadata = {
                "file_type": "html",
                "encoding": self.encoding,
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }
            
            # Try to extract HTML metadata
            try:
                with open(self.file_path, 'r', encoding=self.encoding) as f:
                    html_content = f.read()
                
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract title
                if soup.title:
                    metadata["title"] = soup.title.string.strip()
                
                # Extract meta description
                meta_desc = soup.find("meta", attrs={"name": "description"})
                if meta_desc and meta_desc.get("content"):
                    metadata["description"] = meta_desc["content"]
                
                # Extract meta keywords
                meta_keywords = soup.find("meta", attrs={"name": "keywords"})
                if meta_keywords and meta_keywords.get("content"):
                    metadata["keywords"] = meta_keywords["content"]
                    
            except Exception as e:
                logging.warning(f"Failed to extract HTML metadata: {e}")
            
            return metadata
            
        except ImportError:
            file_path = Path(self.file_path)
            return {
                "file_type": "html",
                "encoding": self.encoding,
                "file_size": file_path.stat().st_size,
                "filename": file_path.name
            }

class SQLProcessor(FileProcessorBase):
    """SQL file processor."""
    
    def __init__(self, file_path: str, encoding: str = "utf-8"):
        self.file_path = file_path
        self.encoding = encoding
    
    def load_documents(self) -> List[str]:
        try:
            with open(self.file_path, 'r', encoding=self.encoding) as f:
                content = f.read().strip()
                
                # Split into statements by semicolon
                statements = [stmt.strip() for stmt in content.split(';') if stmt.strip()]
                
                # Format each statement with a header
                formatted_statements = []
                for i, stmt in enumerate(statements, 1):
                    formatted = f"SQL Statement {i}:\n{stmt}"
                    formatted_statements.append(formatted)
                
                return formatted_statements if formatted_statements else []
                
        except Exception as e:
            logging.error(f"Failed to load SQL file {self.file_path}: {e}")
            raise
    
    def get_metadata(self) -> Dict[str, Any]:
        file_path = Path(self.file_path)
        return {
            "file_type": "sql",
            "encoding": self.encoding,
            "file_size": file_path.stat().st_size,
            "filename": file_path.name
        }

class UniversalFileProcessor:
    """Universal file processor that can handle multiple file types."""
    
    # Supported file types and their extensions
    SUPPORTED_TYPES = {
        '.pdf': PDFProcessor,
        '.txt': TextProcessor,
        '.md': MarkdownProcessor,
        '.markdown': MarkdownProcessor,
        '.csv': CSVProcessor,
        '.docx': DOCXProcessor,
        '.xlsx': ExcelProcessor,
        '.xls': ExcelProcessor,
        '.pptx': PowerPointProcessor,
        '.html': HTMLProcessor,
        '.htm': HTMLProcessor,
        '.sql': SQLProcessor,  # Add SQL support
    }
    
    # MIME type mapping for additional validation
    MIME_TYPES = {
        'application/pdf': ['.pdf'],
        'text/plain': ['.txt', '.sql'],  # Add SQL MIME type
        'text/markdown': ['.md', '.markdown'],
        'text/csv': ['.csv'],
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ['.docx'],
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ['.xlsx'],
        'application/vnd.ms-excel': ['.xls'],
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': ['.pptx'],
        'text/html': ['.html', '.htm'],
    }
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.file_extension = self._get_file_extension()
        self.processor = self._create_processor()
    
    def _get_file_extension(self) -> str:
        """Get the file extension in lowercase."""
        return Path(self.file_path).suffix.lower()
    
    def _create_processor(self) -> FileProcessorBase:
        """Create the appropriate processor for the file type."""
        if self.file_extension not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {self.file_extension}. Supported types: {list(self.SUPPORTED_TYPES.keys())}")
        
        processor_class = self.SUPPORTED_TYPES[self.file_extension]
        return processor_class(self.file_path)
    
    def load_documents(self) -> List[str]:
        """Extract text content from the file."""
        return self.processor.load_documents()
    
    def get_metadata(self) -> Dict[str, Any]:
        """Get metadata from the file."""
        return self.processor.get_metadata()
    
    @classmethod
    def is_supported_file(cls, filename: str) -> bool:
        """Check if a file type is supported."""
        extension = Path(filename).suffix.lower()
        return extension in cls.SUPPORTED_TYPES
    
    @classmethod
    def is_supported_mime_type(cls, mime_type: str) -> bool:
        """Check if a MIME type is supported."""
        return mime_type in cls.MIME_TYPES
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Get list of supported file extensions."""
        return list(cls.SUPPORTED_TYPES.keys())
    
    @classmethod
    def get_supported_mime_types(cls) -> List[str]:
        """Get list of supported MIME types."""
        return list(cls.MIME_TYPES.keys())
    
    @classmethod
    def validate_file(cls, filename: str, mime_type: Optional[str] = None) -> bool:
        """Validate if a file is supported based on filename and/or MIME type."""
        # Check by extension
        if cls.is_supported_file(filename):
            return True
        
        # Check by MIME type if provided
        if mime_type and cls.is_supported_mime_type(mime_type):
            return True
        
        return False 