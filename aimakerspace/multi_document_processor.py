"""
Multi-Document Processor for Regulatory Reporting

This module extends the existing PDF processing capabilities to support
additional document types commonly used in regulatory reporting:
- Excel spreadsheets (COREP, FINREP templates)
- PowerPoint presentations (steering committee decks)
- CSV files (Jira exports, data mappings)
- HTML documents (EBA frameworks)
- Code files (SQL lineage, Python scripts)

The existing PDF processing functionality remains unchanged.
"""

import logging
import os
import tempfile
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
import pandas as pd
from pathlib import Path
import re

# Import existing PDF functionality
try:
    from .pdf_utils import PDFFileLoader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Optional imports for new document types
try:
    import openpyxl
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    import xlrd
    from xlrd import XLRDError
    XLS_AVAILABLE = True
except ImportError:
    XLS_AVAILABLE = False
    xlrd = None

try:
    from pptx import Presentation
    POWERPOINT_AVAILABLE = True
except ImportError:
    POWERPOINT_AVAILABLE = False

try:
    from docx import Document
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False

# Removed HTML processing to simplify dependencies


@dataclass
class ProcessedDocument:
    """Represents a processed document chunk with regulatory metadata"""
    content: str
    metadata: Dict[str, Any]
    doc_type: str
    source_location: str  # e.g., "Page 5", "Sheet1:A1:C10", "Slide 3"
    
    def __post_init__(self):
        """Ensure content is properly formatted"""
        self.content = self.content.strip()


class MultiDocumentProcessor:
    """
    Processor for multiple document types used in regulatory reporting.
    
    This class extends the existing PDF processing capabilities without
    modifying or replacing them.
    """
    
    def __init__(self):
        """Initialize the multi-document processor"""
        self.supported_extensions = {
            '.pdf': 'PDF Documents',
            '.xlsx': 'Excel Spreadsheets',
            '.xls': 'Excel Spreadsheets',
            '.docx': 'Word Documents',
            '.pptx': 'PowerPoint Presentations', 
            '.ppt': 'PowerPoint Presentations',
            '.csv': 'CSV Files',
            '.sql': 'SQL Files',
            '.py': 'Python Files',
            '.js': 'JavaScript Files',
            '.ts': 'TypeScript Files',
            '.md': 'Markdown Files',
            '.txt': 'Text Files'
        }
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions"""
        return list(self.supported_extensions.keys())
    
    def is_supported(self, filename: str) -> bool:
        """Check if a file type is supported"""
        file_ext = Path(filename).suffix.lower()
        return file_ext in self.supported_extensions
    
    def process_document(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """
        Process a document and return structured chunks.
        
        Args:
            file_path: Path to the file
            filename: Original filename
            
        Returns:
            List of ProcessedDocument objects
        """
        file_ext = Path(filename).suffix.lower()
        
        if file_ext == '.pdf':
            return self._process_pdf(file_path, filename)
        elif file_ext in ['.xlsx', '.xls']:
            return self._process_excel(file_path, filename)
        elif file_ext == '.docx':
            return self._process_word(file_path, filename)
        elif file_ext in ['.pptx', '.ppt']:
            return self._process_powerpoint(file_path, filename)
        elif file_ext == '.csv':
            return self._process_csv(file_path, filename)
        elif file_ext in ['.sql', '.py', '.js', '.ts', '.md', '.txt']:
            return self._process_code(file_path, filename)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    def _process_pdf(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process PDF files using existing functionality"""
        if not PDF_AVAILABLE:
            raise ImportError("PDF processing not available")
            
        chunks = []
        try:
            pdf_loader = PDFFileLoader(file_path)
            documents = pdf_loader.load_documents()
            
            for page_num, content in enumerate(documents):
                if content.strip():
                    metadata = {
                        "filename": filename,
                        "page_number": page_num + 1,
                        "doc_type": "pdf",
                        "total_pages": len(documents),
                        "regulatory_type": "basel_document"  # Default assumption for regulatory PDFs
                    }
                    
                    chunk = ProcessedDocument(
                        content=content,
                        metadata=metadata,
                        doc_type="pdf",
                        source_location=f"Page {page_num + 1}"
                    )
                    chunks.append(chunk)
                    
        except Exception as e:
            logging.error(f"Failed to process PDF {filename}: {e}")
            raise
            
        return chunks
    
    def _process_excel(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process Excel files (COREP, FINREP templates)"""
        file_ext = Path(filename).suffix.lower()
        
        if file_ext == '.xlsx':
            return self._process_xlsx(file_path, filename)
        elif file_ext == '.xls':
            return self._process_xls(file_path, filename)
        else:
            raise ValueError(f"Unsupported Excel format: {file_ext}")
    
    def _process_xlsx(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process .xlsx files using openpyxl"""
        if not EXCEL_AVAILABLE:
            raise ImportError("Excel processing requires openpyxl. Install with: pip install openpyxl")
            
        chunks = []
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            
            # Detect regulatory template type
            regulatory_type = self._detect_excel_type(filename, workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                try:
                    worksheet = workbook[sheet_name]
                    sheet_content = self._excel_to_markdown_openpyxl(worksheet, sheet_name)
                    
                    if sheet_content.strip():
                        metadata = {
                            "filename": filename,
                            "sheet_name": sheet_name,
                            "doc_type": "excel",
                            "regulatory_type": regulatory_type,
                            "max_row": worksheet.max_row,
                            "max_column": worksheet.max_column
                        }
                        
                        chunk = ProcessedDocument(
                            content=sheet_content,
                            metadata=metadata,
                            doc_type="excel",
                            source_location=f"Sheet: {sheet_name}"
                        )
                        chunks.append(chunk)
                        
                except Exception as e:
                    logging.warning(f"Failed to process sheet {sheet_name}: {e}")
                    
        except Exception as e:
            logging.error(f"Failed to process Excel file {filename}: {e}")
            raise
            
        return chunks
    
    def _process_xlsx_force(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Force process file with openpyxl regardless of extension"""
        if not EXCEL_AVAILABLE:
            raise ImportError("Excel processing requires openpyxl. Install with: pip install openpyxl")
            
        chunks = []
        try:
            # Force openpyxl to read the file directly
            workbook = load_workbook(file_path, read_only=True, data_only=True)
            
            # Detect regulatory template type
            regulatory_type = self._detect_excel_type(filename, workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                try:
                    worksheet = workbook[sheet_name]
                    sheet_content = self._excel_to_markdown_openpyxl(worksheet, sheet_name)
                    
                    if sheet_content.strip():
                        metadata = {
                            "filename": filename,
                            "sheet_name": sheet_name,
                            "doc_type": "excel",
                            "regulatory_type": regulatory_type,
                            "max_row": worksheet.max_row,
                            "max_column": worksheet.max_column,
                            "format_note": "processed_as_xlsx_despite_xls_extension"
                        }
                        
                        chunk = ProcessedDocument(
                            content=sheet_content,
                            metadata=metadata,
                            doc_type="excel",
                            source_location=f"Sheet: {sheet_name}"
                        )
                        chunks.append(chunk)
                        
                except Exception as e:
                    logging.warning(f"Failed to process sheet {sheet_name}: {e}")
                    
        except Exception as e:
            logging.error(f"Failed to force process Excel file {filename} with openpyxl: {e}")
            raise
            
        return chunks
    
    def _process_xls(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process .xls files using xlrd"""
        if not XLS_AVAILABLE:
            raise ImportError("Legacy Excel processing requires xlrd. Install with: pip install xlrd")
            
        chunks = []
        try:
            workbook = xlrd.open_workbook(file_path)
            sheet_names = workbook.sheet_names()
            
            # Detect regulatory template type
            regulatory_type = self._detect_excel_type(filename, sheet_names)
            
            for sheet_name in sheet_names:
                try:
                    worksheet = workbook.sheet_by_name(sheet_name)
                    sheet_content = self._excel_to_markdown_xlrd(worksheet, sheet_name)
                    
                    if sheet_content.strip():
                        metadata = {
                            "filename": filename,
                            "sheet_name": sheet_name,
                            "doc_type": "excel",
                            "regulatory_type": regulatory_type,
                            "max_row": worksheet.nrows,
                            "max_column": worksheet.ncols
                        }
                        
                        chunk = ProcessedDocument(
                            content=sheet_content,
                            metadata=metadata,
                            doc_type="excel",
                            source_location=f"Sheet: {sheet_name}"
                        )
                        chunks.append(chunk)
                        
                except Exception as e:
                    logging.warning(f"Failed to process sheet {sheet_name}: {e}")
                    continue
                    
            # xlrd Book objects don't need explicit closing
            
        except Exception as e:
            error_msg = str(e).lower()
            if ("xlsx" in error_msg or "not supported" in error_msg or 
                "xlrdError" in str(type(e).__name__)):
                # File is actually .xlsx format, try with openpyxl by forcing it
                logging.info(f"File {filename} appears to be .xlsx format, trying openpyxl directly")
                try:
                    return self._process_xlsx_force(file_path, filename)
                except Exception as fallback_error:
                    logging.error(f"Failed to process {filename} with both xlrd and openpyxl: {e}, {fallback_error}")
                    raise
            else:
                logging.error(f"Failed to process Excel file {filename}: {e}")
                raise
            
        return chunks
    
    def _detect_excel_type(self, filename: str, sheet_names: List[str]) -> str:
        """Detect the type of regulatory Excel template"""
        filename_lower = filename.lower()
        sheet_names_str = " ".join(sheet_names).lower()
        
        if any(term in filename_lower for term in ['corep', 'capital', 'own funds']):
            return "corep_template"
        elif any(term in filename_lower for term in ['finrep', 'financial', 'ifrs']):
            return "finrep_template"
        elif any(term in filename_lower for term in ['mapping', 'lineage', 'source']):
            return "data_mapping"
        elif any(term in sheet_names_str for term in ['corep', 'capital']):
            return "corep_template"
        elif any(term in sheet_names_str for term in ['finrep', 'financial']):
            return "finrep_template"
        else:
            return "regulatory_template"
    
    def _excel_to_markdown_openpyxl(self, worksheet, sheet_name: str) -> str:
        """Convert Excel sheet to markdown format"""
        content = [f"# Regulatory Template: {sheet_name}\n"]
        
        # Find actual data range
        max_row = min(worksheet.max_row or 0, 50)  # Limit to 50 rows
        max_col = min(worksheet.max_column or 0, 10)  # Limit to 10 columns
        
        if max_row == 0 or max_col == 0:
            return f"# Sheet: {sheet_name}\n\n*Empty sheet*"
        
        # Convert to table
        rows = []
        for row_num in range(1, max_row + 1):
            row_data = []
            for col_num in range(1, max_col + 1):
                cell = worksheet.cell(row=row_num, column=col_num)
                value = str(cell.value) if cell.value is not None else ""
                value = value.replace("|", "\\|").replace("\n", " ").strip()
                row_data.append(value)
            rows.append(row_data)
        
        if rows:
            content.append("| " + " | ".join(rows[0]) + " |")
            content.append("|" + "|".join([" --- " for _ in rows[0]]) + "|")
            
            for row in rows[1:]:
                content.append("| " + " | ".join(row) + " |")
        
        return "\n".join(content)
    
    def _excel_to_markdown_xlrd(self, worksheet, sheet_name: str) -> str:
        """Convert Excel sheet to markdown format using xlrd"""
        content = [f"# Regulatory Template: {sheet_name}\n"]
        
        # Process rows and columns
        for row_idx in range(min(50, worksheet.nrows)):  # Limit to first 50 rows
            row_data = []
            has_content = False
            
            for col_idx in range(min(10, worksheet.ncols)):  # Limit to first 10 columns
                try:
                    cell_value = worksheet.cell_value(row_idx, col_idx)
                    if cell_value:
                        has_content = True
                        # Handle different cell types
                        if isinstance(cell_value, float) and cell_value.is_integer():
                            cell_value = int(cell_value)
                        row_data.append(str(cell_value))
                    else:
                        row_data.append("")
                except Exception:
                    row_data.append("")
            
            if has_content:
                if row_idx == 0:
                    # Header row
                    content.append("| " + " | ".join(row_data) + " |")
                    content.append("| " + " | ".join(["---"] * len(row_data)) + " |")
                else:
                    content.append("| " + " | ".join(row_data) + " |")
        
        return "\n".join(content)
    
    def _process_powerpoint(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process PowerPoint presentations (steering committee decks)"""
        if not POWERPOINT_AVAILABLE:
            raise ImportError("PowerPoint processing requires python-pptx. Install with: pip install python-pptx")
            
        chunks = []
        try:
            presentation = Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides):
                try:
                    slide_content = self._extract_slide_content(slide, slide_num + 1)
                    
                    if slide_content.strip():
                        metadata = {
                            "filename": filename,
                            "slide_number": slide_num + 1,
                            "doc_type": "powerpoint",
                            "regulatory_type": "steering_committee",
                            "total_slides": len(presentation.slides)
                        }
                        
                        chunk = ProcessedDocument(
                            content=slide_content,
                            metadata=metadata,
                            doc_type="powerpoint",
                            source_location=f"Slide {slide_num + 1}"
                        )
                        chunks.append(chunk)
                        
                except Exception as e:
                    logging.warning(f"Failed to extract content from slide {slide_num + 1}: {e}")
                    continue
                    
        except Exception as e:
            logging.error(f"Failed to process PowerPoint file {filename}: {e}")
            raise
            
        return chunks
    
    def _extract_slide_content(self, slide, slide_num: int) -> str:
        """Extract content from PowerPoint slide"""
        content = [f"# Regulatory Presentation - Slide {slide_num}\n"]
        
        # Extract title
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
            if title:
                content.append(f"## {title}\n")
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                content.append(text)
                content.append("")
        
        return "\n".join(content)
    
    def _process_csv(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process CSV files (Jira exports, data mappings)"""
        chunks = []
        try:
            # Try different encodings and error handling strategies
            df = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                for error_bad_lines in [False, True]:  # First try skipping bad lines, then be strict
                    for quotechar in ['"', "'"]:  # Try different quote characters
                        try:
                            df = pd.read_csv(
                                file_path, 
                                encoding=encoding,
                                on_bad_lines='skip' if error_bad_lines else 'error',
                                quotechar=quotechar,
                                engine='python'  # More flexible parser
                            )
                            logging.info(f"Successfully parsed CSV with encoding={encoding}, skip_bad_lines={error_bad_lines}, quotechar={quotechar}")
                            break
                        except (UnicodeDecodeError, pd.errors.ParserError):
                            continue
                    if df is not None:
                        break
                if df is not None:
                    break
            
            if df is None:
                # Final fallback: read as plain text and create simple summary
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                
                summary_content = f"# Regulatory CSV: {filename}\n\n"
                summary_content += "**Note:** File had parsing issues, showing raw preview\n\n"
                summary_content += "## File Preview:\n"
                summary_content += "```\n"
                summary_content += content[:2000]  # First 2000 characters
                if len(content) > 2000:
                    summary_content += "\n... (truncated)"
                summary_content += "\n```"
                
                metadata = {
                    "filename": filename,
                    "doc_type": "csv",
                    "regulatory_type": "unknown",
                    "total_rows": "unknown",
                    "columns": [],
                    "column_count": 0,
                    "parsing_status": "failed"
                }
                
                chunk = ProcessedDocument(
                    content=summary_content,
                    metadata=metadata,
                    doc_type="csv",
                    source_location="CSV Raw Content (parsing failed)"
                )
                chunks.append(chunk)
                return chunks
            
            # Detect CSV type
            csv_type = self._detect_csv_type(filename, df.columns.tolist())
            
            # Create summary
            summary_content = self._create_csv_summary(df, filename, csv_type)
            
            metadata = {
                "filename": filename,
                "doc_type": "csv",
                "regulatory_type": csv_type,
                "total_rows": len(df),
                "columns": list(df.columns),
                "column_count": len(df.columns),
                "parsing_status": "success"
            }
            
            chunk = ProcessedDocument(
                content=summary_content,
                metadata=metadata,
                doc_type="csv",
                source_location=f"CSV Summary ({len(df)} rows)"
            )
            chunks.append(chunk)
            
        except Exception as e:
            logging.error(f"Failed to process CSV file {filename}: {e}")
            raise
            
        return chunks
    
    def _detect_csv_type(self, filename: str, columns: List[str]) -> str:
        """Detect the type of CSV file"""
        filename_lower = filename.lower()
        columns_str = " ".join(columns).lower()
        
        if any(term in filename_lower for term in ['jira', 'issue', 'ticket']):
            return "jira_export"
        elif any(term in columns_str for term in ['issue', 'key', 'status', 'assignee']):
            return "jira_export"
        elif any(term in filename_lower for term in ['mapping', 'lineage', 'source']):
            return "data_mapping"
        else:
            return "regulatory_data"
    
    def _create_csv_summary(self, df: pd.DataFrame, filename: str, csv_type: str) -> str:
        """Create a summary of the CSV file"""
        content = [f"# Regulatory CSV: {filename}\n"]
        content.append(f"**Type:** {csv_type}")
        content.append(f"**Total Rows:** {len(df)}")
        content.append(f"**Total Columns:** {len(df.columns)}\n")
        
        content.append("## Columns:")
        for col in df.columns:
            content.append(f"- {col}")
        
        content.append("\n## Sample Data:")
        sample_df = df.head()
        content.append(sample_df.to_string())
        
        return "\n".join(content)
    
    def _process_word(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process Word documents (regulatory policies, procedures)"""
        if not WORD_AVAILABLE:
            raise ImportError("Word processing requires python-docx. Install with: pip install python-docx")
            
        chunks = []
        try:
            doc = Document(file_path)
            
            # Collect all text content
            all_paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    all_paragraphs.append(para.text.strip())
            
            # Detect document type
            full_text = " ".join(all_paragraphs)
            word_doc_type = self._detect_word_doc_type(filename, full_text)
            
            # Create chunks of reasonable size (combine paragraphs)
            current_chunk = []
            current_size = 0
            chunk_limit = 1000  # characters
            
            for para_text in all_paragraphs:
                if current_size + len(para_text) > chunk_limit and current_chunk:
                    # Create chunk from current paragraphs
                    chunk_content = "\n\n".join(current_chunk)
                    
                    metadata = {
                        "filename": filename,
                        "doc_type": "word",
                        "regulatory_type": word_doc_type,
                        "total_paragraphs": len(all_paragraphs),
                        "chunk_size": len(chunk_content)
                    }
                    
                    chunk = ProcessedDocument(
                        content=chunk_content,
                        metadata=metadata,
                        doc_type="word",
                        source_location=f"Document Section ({len(current_chunk)} paragraphs)"
                    )
                    chunks.append(chunk)
                    
                    # Start new chunk
                    current_chunk = [para_text]
                    current_size = len(para_text)
                else:
                    current_chunk.append(para_text)
                    current_size += len(para_text)
            
            # Add final chunk if any content remains
            if current_chunk:
                chunk_content = "\n\n".join(current_chunk)
                
                metadata = {
                    "filename": filename,
                    "doc_type": "word",
                    "regulatory_type": word_doc_type,
                    "total_paragraphs": len(all_paragraphs),
                    "chunk_size": len(chunk_content)
                }
                
                chunk = ProcessedDocument(
                    content=chunk_content,
                    metadata=metadata,
                    doc_type="word",
                    source_location=f"Document Section ({len(current_chunk)} paragraphs)"
                )
                chunks.append(chunk)
                    
        except Exception as e:
            logging.error(f"Failed to process Word file {filename}: {e}")
            raise
            
        return chunks
    
    def _detect_word_doc_type(self, filename: str, content: str) -> str:
        """Detect the type of Word document for regulatory context"""
        filename_lower = filename.lower()
        content_lower = content.lower()
        
        if any(term in filename_lower for term in ['policy', 'procedure', 'manual']):
            return "regulatory_policy"
        elif any(term in filename_lower for term in ['corep', 'finrep', 'basel']):
            return "regulatory_guidance"
        elif any(term in content_lower for term in ['policy', 'procedure', 'shall', 'must']):
            return "regulatory_policy"
        elif any(term in content_lower for term in ['capital', 'liquidity', 'risk management']):
            return "regulatory_guidance"
        else:
            return "regulatory_document"
    
    def _process_code(self, file_path: str, filename: str) -> List[ProcessedDocument]:
        """Process code files (SQL lineage, Python scripts)"""
        chunks = []
        try:
            file_ext = Path(filename).suffix.lower()
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            language = self._detect_language(file_ext)
            code_type = self._detect_code_type(filename, content)
            
            line_count = len(content.split('\n'))
            formatted_content = [
                f"# Regulatory Code: {filename}\n",
                f"**Language:** {language}",
                f"**Type:** {code_type}",
                f"**Lines:** {line_count}\n",
                f"```{language}",
                content,
                "```"
            ]
            
            metadata = {
                "filename": filename,
                "doc_type": "code",
                "regulatory_type": code_type,
                "language": language,
                "file_extension": file_ext,
                "line_count": line_count
            }
            
            chunk = ProcessedDocument(
                content="\n".join(formatted_content),
                metadata=metadata,
                doc_type="code",
                source_location="Full File"
            )
            chunks.append(chunk)
            
        except Exception as e:
            logging.error(f"Failed to process code file {filename}: {e}")
            raise
            
        return chunks
    
    def _detect_language(self, file_ext: str) -> str:
        """Detect programming language from file extension"""
        language_map = {
            '.sql': 'sql',
            '.py': 'python',
            '.js': 'javascript',
            '.ts': 'typescript',
            '.md': 'markdown',
            '.txt': 'text'
        }
        return language_map.get(file_ext, 'text')
    
    def _detect_code_type(self, filename: str, content: str) -> str:
        """Detect the type of code file for regulatory context"""
        filename_lower = filename.lower()
        content_lower = content.lower()
        
        if any(term in filename_lower for term in ['lineage', 'etl', 'mapping']):
            return "data_lineage"
        elif any(term in filename_lower for term in ['corep', 'finrep', 'basel']):
            return "regulatory_calculation"
        elif any(term in content_lower for term in ['select', 'from', 'join', 'where']):
            return "sql_query"
        else:
            return "regulatory_script" 